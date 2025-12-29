import os
import argparse
from pathlib import Path
import docx
import openpyxl
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import datetime
from markitdown import MarkItDown
import zipfile
import shutil
import tempfile
import re
import subprocess
import chardet
import tarfile
import logging
import json
from typing import Optional, Dict, List, Any
from dataclasses import dataclass, field, asdict

try:
    import py7zr
    HAS_7Z = True
except ImportError:
    HAS_7Z = False

try:
    import rarfile
    HAS_RAR = True
except ImportError:
    HAS_RAR = False

try:
    import lhafile
    HAS_LZH = True
except ImportError:
    HAS_LZH = False

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import magic
    HAS_MAGIC = True
    # MIMEタイプ判定用インスタンス（再利用）
    _mime_detector = magic.Magic(mime=True)
except ImportError:
    HAS_MAGIC = False
    _mime_detector = None


# ---------------------------------------------------------
# Logging Setup
# ---------------------------------------------------------

def setup_logging(output_dir: Path, verbose: bool = False) -> logging.Logger:
    """
    ログ機構をセットアップする
    
    Args:
        output_dir: 出力ディレクトリ（ログファイル出力先）
        verbose: 詳細ログ出力フラグ
        
    Returns:
        設定済みのロガー
    """
    log_level = logging.DEBUG if verbose else logging.INFO
    
    # ログディレクトリ作成
    log_dir = output_dir / "logs"
    log_dir.mkdir(parents=True, exist_ok=True)
    
    # ログファイル名（タイムスタンプ付き）
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    log_file = log_dir / f"processing_{timestamp}.log"
    
    # ロガー設定
    logger = logging.getLogger("notebooklm_loader")
    logger.setLevel(log_level)
    
    # 既存ハンドラをクリア
    logger.handlers.clear()
    
    # ファイルハンドラ
    file_handler = logging.FileHandler(log_file, encoding='utf-8')
    file_handler.setLevel(logging.DEBUG)
    file_format = logging.Formatter('%(asctime)s | %(levelname)-8s | %(message)s')
    file_handler.setFormatter(file_format)
    logger.addHandler(file_handler)
    
    # コンソールハンドラ
    console_handler = logging.StreamHandler()
    console_handler.setLevel(log_level)
    console_format = logging.Formatter('%(message)s')
    console_handler.setFormatter(console_format)
    logger.addHandler(console_handler)
    
    logger.info(f"Log file: {log_file}")
    return logger


# ---------------------------------------------------------
# Processing Summary
# ---------------------------------------------------------

@dataclass
class FileResult:
    """個別ファイルの処理結果"""
    path: str
    status: str  # converted, skipped, error, password_protected
    output: Optional[str] = None
    error_message: Optional[str] = None
    file_type: Optional[str] = None


@dataclass
class ProcessingSummary:
    """処理サマリー"""
    run_time: str = field(default_factory=lambda: datetime.datetime.now().isoformat())
    target_path: str = ""
    total_files: int = 0
    processed: int = 0
    skipped: int = 0
    errors: int = 0
    password_protected: int = 0
    files: List[Dict[str, Any]] = field(default_factory=list)
    
    def add_result(self, result: FileResult):
        """処理結果を追加"""
        self.files.append(asdict(result))
        self.total_files += 1
        
        if result.status == "converted":
            self.processed += 1
        elif result.status == "skipped":
            self.skipped += 1
        elif result.status == "error":
            self.errors += 1
        elif result.status == "password_protected":
            self.password_protected += 1
    
    def save(self, output_dir: Path):
        """サマリーをJSONファイルに保存"""
        summary_file = output_dir / "processing_report.json"
        with open(summary_file, 'w', encoding='utf-8') as f:
            json.dump(asdict(self), f, ensure_ascii=False, indent=2)
        return summary_file

# ---------------------------------------------------------
# PDF Conversion Utilities
# ---------------------------------------------------------

def convert_to_pdf_via_libreoffice(input_path, output_dir_path):
    """
    LibreOffice (soffice) を使用してPDF変換を行う。
    Mac: /Applications/LibreOffice.app/Contents/MacOS/soffice
    """
    soffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if not os.path.exists(soffice_path):
        soffice_path = "soffice" # Try PATH

    cmd = [
        soffice_path,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(output_dir_path),
        str(input_path)
    ]
    try:
        # Run command
        subprocess.run(cmd, check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        
        # Output filename check (LibreOffice creates filename.pdf)
        original_stem = input_path.stem
        
        generated_pdf = output_dir_path / (original_stem + ".pdf")
        if generated_pdf.exists():
            return generated_pdf
        return None
    except Exception as e:
        print(f"    [PDF Convert Error] {e}")
        return None

def convert_image_to_pdf(input_path, output_dir_path):
    """
    画像ファイルをPDFに変換する（Pillow使用）
    """
    if not HAS_PIL:
        print(f"    [Warning] Pillow not installed, skipping image: {input_path.name}")
        return None
    
    try:
        img = Image.open(input_path)
        # RGBAの場合はRGBに変換（PDF保存のため）
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')
        
        output_pdf = output_dir_path / (input_path.stem + ".pdf")
        img.save(output_pdf, "PDF", resolution=100.0)
        return output_pdf
    except Exception as e:
        print(f"    [Image to PDF Error] {e}")
        return None

# ---------------------------------------------------------
# 定数定義
# ---------------------------------------------------------

OUTPUT_DIR_NAME = "converted_files"
COMBINED_FILENAME = "All_Files_Combined.txt"

# 判定閾値（目安）
TEXT_PER_VISUAL_THRESHOLD = 300 

# 巨大ファイルスキップ閾値
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB

# Office拡張子（新形式 - 視覚密度分析可能）
OFFICE_EXTENSIONS_NEW = {'.docx', '.xlsx', '.pptx'}

# Office拡張子（旧形式 - MarkItDownで直接変換）
OFFICE_EXTENSIONS_LEGACY = {'.doc', '.xls', '.ppt'}

# 全Office拡張子
OFFICE_EXTENSIONS_ALL = OFFICE_EXTENSIONS_NEW | OFFICE_EXTENSIONS_LEGACY

# MarkItDownで直接変換可能な形式
MARKITDOWN_EXTENSIONS = {
    '.rtf',    # リッチテキスト
    '.epub',   # 電子書籍
    '.msg',    # Outlook メール
    '.eml',    # メール
}

# Visio（PDF変換対象）
VISIO_EXTENSIONS = {'.vsdx', '.vsd'}

# 画像（PDF変換対象）
IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.webp'}

# 圧縮形式
ARCHIVE_EXTENSIONS = {'.zip', '.7z', '.rar', '.tar', '.gz', '.tgz', '.lzh'}

# スキップ対象（処理不可または対象外）
SKIP_EXTENSIONS = {
    # OneNote
    '.one', '.onetoc2',
    # Access
    '.accdb', '.mdb',
    # 動画
    '.mp4', '.avi', '.mov', '.wmv', '.mkv', '.flv', '.webm',
    # 音声
    '.mp3', '.wav', '.aac', '.flac', '.ogg', '.wma', '.m4a',
    # CAD
    '.dwg', '.dxf',
    # 実行ファイル
    '.exe', '.dll', '.so', '.dylib',
    # その他バイナリ
    '.bin', '.dat', '.iso', '.img',
}

TEXT_EXTENSIONS = {
    '.txt', '.md', '.py', '.js', '.jsx', '.ts', '.tsx', '.html', '.css', '.json', 
    '.yaml', '.yml', '.org', '.sh', '.bat', '.zsh', '.rb', '.java', '.c', '.cpp', 
    '.h', '.go', '.rs', '.php', '.pl', '.swift', '.kt', '.sql', '.xml', '.csv',
    '.log', '.ini', '.cfg', '.conf', '.properties', '.env', '.toml', '.tsv', '.rst'
}

def setup_args():
    """
    コマンドライン引数をパースする
    
    Returns:
        パース済みの引数オブジェクト
    """
    parser = argparse.ArgumentParser(
        description='Office files to Markdown converter for NotebookLM',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  %(prog)s /path/to/folder                    # 基本変換
  %(prog)s /path/to/folder --merge            # スマート結合モード
  %(prog)s /path/to/folder --merge --verbose  # 詳細ログ出力
  %(prog)s /path/to/folder --dry-run          # 実行計画のみ表示
        """
    )
    parser.add_argument('target_dir', help='Target directory containing Office files or ZIP file')
    parser.add_argument('--merge', action='store_true', 
                        help='Also create merged output in converted_files_merged directory')
    parser.add_argument('--skip-ppt', action='store_true', 
                        help='Skip PowerPoint files (recommend using PDF for visual-heavy PPTs)')
    parser.add_argument('-v', '--verbose', action='store_true',
                        help='Enable verbose logging (DEBUG level)')
    parser.add_argument('-q', '--quiet', action='store_true',
                        help='Suppress console output (only log to file)')
    parser.add_argument('--dry-run', action='store_true',
                        help='Show what would be processed without actually converting')
    return parser.parse_args()

def analyze_docx(file_path):
    """
    Wordファイルの視覚要素と文字数を分析する
    
    Args:
        file_path: 対象ファイルのパス
        
    Returns:
        (visual_count, char_count): 視覚要素数と文字数のタプル
    """
    logger = logging.getLogger("notebooklm_loader")
    try:
        doc = docx.Document(file_path)
        visual_count = 0
        char_count = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            char_count += len(text)
            for run in para.runs:
                if run.element.xpath('.//a:blip'):
                     visual_count += 1
        return visual_count, char_count
    except Exception as e:
        logger.debug(f"analyze_docx error for {file_path}: {e}")
        return 0, 0

def analyze_xlsx(file_path):
    """
    Excelファイルの視覚要素と文字数を分析する
    
    Args:
        file_path: 対象ファイルのパス
        
    Returns:
        (visual_count, char_count): 視覚要素数と文字数のタプル
    """
    logger = logging.getLogger("notebooklm_loader")
    try:
        visual_count = 0
        char_count = 0
        try:
            wb_obj = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb_obj.sheetnames:
                sheet = wb_obj[sheet_name]
                if hasattr(sheet, '_charts') and sheet._charts:
                    visual_count += len(sheet._charts)
        except Exception as e:
            logger.debug(f"analyze_xlsx chart detection error: {e}")
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            try:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                if not df.empty:
                    csv_text = df.to_csv(index=False)
                    char_count += len(csv_text)
            except Exception as e:
                logger.debug(f"analyze_xlsx sheet {sheet_name} error: {e}")
        return visual_count, char_count
    except Exception as e:
        logger.debug(f"analyze_xlsx error for {file_path}: {e}")
        return 0, 0

def analyze_pptx(file_path):
    """
    PowerPointファイルの視覚要素と文字数を分析する
    
    Args:
        file_path: 対象ファイルのパス
        
    Returns:
        (visual_count, char_count): 視覚要素数と文字数のタプル
    """
    logger = logging.getLogger("notebooklm_loader")
    try:
        prs = Presentation(file_path)
        visual_count = 0
        char_count = 0
        for slide in prs.slides:
            if slide.shapes.title:
                char_count += len(slide.shapes.title.text.strip())
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    char_count += len(shape.text_frame.text.strip())
                is_visual = False
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: is_visual = True
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP: is_visual = True
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                     if not shape.has_text_frame or not shape.text_frame.text.strip(): is_visual = True
                if is_visual: visual_count += 1
        return visual_count, char_count
    except Exception as e:
        logger.debug(f"analyze_pptx error for {file_path}: {e}")
        return 0, 0

def convert_with_markitdown(file_path):
    try:
        md = MarkItDown()
        result = md.convert(str(file_path)) 
        if result and result.text_content:
            return result.text_content
        return ""
    except Exception as e:
        print(f"    Error converting {file_path.name}: {e}")
        return None

def sanitize_filename(name):
    """ファイル名として使える文字だけに置換"""
    return re.sub(r'[\\/*?:"<>|]', "", name)

def get_output_filename(root_path, file_path, extension=".md"):
    """
    元のフォルダ構造を反映したファイル名を生成する。
    例: root/A/B/file.docx -> A_B_file.md
    """
    try:
        rel_path = file_path.relative_to(root_path)
        # フォルダ区切りをアンダースコア等に置換してフラット化
        flat_name = str(rel_path.with_suffix('')).replace(os.sep, '_')
        return sanitize_filename(flat_name) + extension
    except ValueError:
        # root_path外にある場合（稀だが一応）
        return file_path.stem + extension


def extract_zip_with_encoding(zip_path, extract_to):
    """
    ZIPファイルを解凍する際、Windows等で作成されたShift-JISのファイル名を
    正しく復元して展開する。パスワード保護されている場合は警告を返す。
    """
    try:
        with zipfile.ZipFile(zip_path, 'r') as z:
            # パスワード保護チェック
            for file_info in z.infolist():
                if file_info.flag_bits & 0x1:  # 暗号化フラグ
                    return "PASSWORD_PROTECTED"
            
            for file_info in z.infolist():
                filename = file_info.filename
                
                # UTF-8フラグが立っていない場合、エンコーディングの補正を試みる
                if file_info.flag_bits & 0x800 == 0:
                    try:
                        # Windows (Japanese) ZIP is often CP932 encoded but marked as CP437
                        filename = filename.encode('cp437').decode('cp932')
                    except (UnicodeDecodeError, UnicodeEncodeError):
                        pass  # エンコーディング変換失敗は元のファイル名を使用
                
                # ターゲットパスの生成
                target_path = Path(extract_to) / filename
                
                # ディレクトリトラバーサル対策
                if not os.path.abspath(target_path).startswith(os.path.abspath(extract_to)):
                    continue
                    
                if file_info.is_dir():
                    target_path.mkdir(parents=True, exist_ok=True)
                else:
                    target_path.parent.mkdir(parents=True, exist_ok=True)
                    with z.open(file_info) as source, open(target_path, "wb") as target:
                        shutil.copyfileobj(source, target)
        return "OK"
    except RuntimeError as e:
        if "password" in str(e).lower() or "encrypted" in str(e).lower():
            return "PASSWORD_PROTECTED"
        raise


def extract_7z(archive_path, extract_to):
    """7zファイルを展開"""
    if not HAS_7Z:
        print(f"    [Warning] py7zr not installed, skipping: {archive_path.name}")
        return "LIBRARY_MISSING"
    try:
        with py7zr.SevenZipFile(archive_path, mode='r') as z:
            if z.needs_password():
                return "PASSWORD_PROTECTED"
            z.extractall(path=extract_to)
        return "OK"
    except Exception as e:
        if "password" in str(e).lower():
            return "PASSWORD_PROTECTED"
        print(f"    [7z Extract Error] {e}")
        return "ERROR"


def extract_rar(archive_path, extract_to):
    """RARファイルを展開"""
    if not HAS_RAR:
        print(f"    [Warning] rarfile not installed, skipping: {archive_path.name}")
        return "LIBRARY_MISSING"
    try:
        with rarfile.RarFile(archive_path) as rf:
            if rf.needs_password():
                return "PASSWORD_PROTECTED"
            rf.extractall(path=extract_to)
        return "OK"
    except rarfile.NeedFirstVolume:
        print(f"    [Warning] Multi-volume RAR, skipping: {archive_path.name}")
        return "MULTI_VOLUME"
    except Exception as e:
        if "password" in str(e).lower():
            return "PASSWORD_PROTECTED"
        print(f"    [RAR Extract Error] {e}")
        return "ERROR"


def extract_tar(archive_path, extract_to):
    """
    tar/tar.gz/tgzファイルを展開
    
    Args:
        archive_path: 圧縮ファイルのパス
        extract_to: 展開先ディレクトリ
        
    Returns:
        処理結果（"OK", "ERROR"等）
    """
    import sys
    try:
        with tarfile.open(archive_path, 'r:*') as tf:
            # ディレクトリトラバーサル対策
            for member in tf.getmembers():
                member_path = os.path.join(extract_to, member.name)
                if not os.path.abspath(member_path).startswith(os.path.abspath(extract_to)):
                    continue
                # Python 3.12+ではfilter引数が必要
                if sys.version_info >= (3, 12):
                    tf.extract(member, extract_to, filter='data')
                else:
                    tf.extract(member, extract_to)
        return "OK"
    except Exception as e:
        print(f"    [TAR Extract Error] {e}")
        return "ERROR"


def extract_lzh(archive_path, extract_to):
    """LZHファイルを展開"""
    if not HAS_LZH:
        print(f"    [Warning] lhafile not installed, skipping: {archive_path.name}")
        return "LIBRARY_MISSING"
    try:
        with lhafile.LhaFile(str(archive_path)) as lf:
            for info in lf.infolist():
                target_path = Path(extract_to) / info.filename
                # ディレクトリトラバーサル対策
                if not os.path.abspath(target_path).startswith(os.path.abspath(extract_to)):
                    continue
                target_path.parent.mkdir(parents=True, exist_ok=True)
                with open(target_path, 'wb') as f:
                    f.write(lf.read(info.filename))
        return "OK"
    except Exception as e:
        print(f"    [LZH Extract Error] {e}")
        return "ERROR"


def get_mime_type(file_path):
    """
    ファイルのMIMEタイプを取得する（python-magic使用）
    
    Args:
        file_path: 対象ファイルのパス
        
    Returns:
        MIMEタイプ文字列、またはNone（ライブラリ未インストール時）
    """
    if not _mime_detector:
        return None
    try:
        return _mime_detector.from_file(str(file_path))
    except Exception:
        return None


def is_likely_text_by_mime(file_path):
    """MIMEタイプからテキストファイルかどうか判定"""
    mime = get_mime_type(file_path)
    if mime is None:
        return None  # 判定不可
    
    text_mimes = [
        'text/', 'application/json', 'application/xml',
        'application/javascript', 'application/x-sh',
    ]
    return any(mime.startswith(t) for t in text_mimes)



# ---------------------------------------------------------
# Feature: Smart Chunking & Merged Output
# ---------------------------------------------------------

MAX_CHARS_PER_VOLUME = 10500000  # Approx 35MB limit (Safe margin for NotebookLM's 40MB limit)

class MergedOutputManager:
    def __init__(self, output_dir):
        self.output_dir = output_dir
        self.output_dir.mkdir(exist_ok=True)
        self.current_vol = 1
        self.current_content = []
        self.current_char_count = 0
        self.file_index = [] # List of filenames included in current vol

    def add_content(self, filename, content):
        """
        コンテンツを追加する。
        もしコンテンツ単体で制限を超える場合は、分割して追加する（Recursive Split）。
        追加によって制限を超える場合は、現在のVolを書き出して次へ行く。
        """
        content_len = len(content)

        # Case 1: Huge single file -> Recursive Split
        # (ヘッダ込みで計算済みだが、念のためここでもチェックする論理は変えないが、handle_huge_file内で厳密計算する)
        if content_len > MAX_CHARS_PER_VOLUME:
            self._handle_huge_file(filename, content)
            return

        # Case 2: Buffer overflow -> Flush and define new volume
        if self.current_char_count + content_len > MAX_CHARS_PER_VOLUME:
            self._flush_volume()

        # Normal Add
        self.current_content.append(content)
        self.current_char_count += content_len
        self.file_index.append(filename)

    def _handle_huge_file(self, filename, content):
        """巨大ファイルを分割して登録する"""
        # content は既にヘッダーがついている状態だが、
        # ここでは再分割するため、ヘッダーを除去して本文だけ取り出したいところだが、
        # 引数の content は 'final_content' であり、メタデータを含んでいる。
        # 厳密にはメタデータごと分割するのは変なので、
        # 本文だけ抽出するよりは、渡される前の raw content を引数に取るべきだが、
        # 設計上、Markdown変換後の最終テキストを受け取っているので、
        # ここでは「巨大なテキスト塊」として扱い、強制分割する方針とする。
        
        remaining = content
        part_num = 1
        
        while remaining:
            # 次のパート用ヘッダー（概算サイズ）
            part_header = f"\n\n# {filename} (Part {part_num})\n\n"
            header_len = len(part_header)
            
            # 残りの容量
            available_space = MAX_CHARS_PER_VOLUME - header_len
            
            # 現在のバッファに空きがなければフラッシュ
            if self.current_char_count + header_len > MAX_CHARS_PER_VOLUME: # ヘッダすら入らないならフラッシュ
                 self._flush_volume()
            
            # 再計算（フラッシュ後）
            available_space = MAX_CHARS_PER_VOLUME - self.current_char_count - header_len
            
            if len(remaining) > available_space:
                # カット
                c_chunk = remaining[:available_space]
                remaining = remaining[available_space:]
            else:
                c_chunk = remaining
                remaining = ""
            
            full_chunk = part_header + c_chunk
            
            # ここまできたら必ず入るはず
            self.current_content.append(full_chunk)
            self.file_index.append(f"{filename} (Part {part_num})")
            self.current_char_count += len(full_chunk)
            
            # 満杯になったらフラッシュ
            if self.current_char_count >= MAX_CHARS_PER_VOLUME:
                self._flush_volume()
            
            part_num += 1

    def _flush_volume(self):
        """現在のバッファをファイルに書き出す"""
        if not self.current_content:
            return

        if self.current_content: # Double check
            vol_filename = f"Merged_Files_Vol{self.current_vol:02d}.md"
            output_path = self.output_dir / vol_filename
            
            # 目次生成
            index_text = "# Table of Contents\n" + "\n".join([f"- {name}" for name in self.file_index]) + "\n\n---\n\n"
            
            full_text = index_text + "\n".join(self.current_content)
            
            try:
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(full_text)
                print(f"[Merged Created] {vol_filename} ({len(full_text)} chars)")
            except Exception as e:
                print(f"Error writing volume {vol_filename}: {e}")

            # Reset
            self.current_vol += 1
            self.current_content = []
            self.current_char_count = 0
            self.file_index = []

    def finalize(self):
        """最後に残っているバッファを書き出す"""
        self._flush_volume()


def is_text_file(file_path):
    """
    chardetを使ってテキストファイルかどうか判定する。
    テキストの場合は (True, encoding) を返す。
    バイナリの場合は (False, None) を返す。
    """
    try:
        with open(file_path, 'rb') as f:
            raw = f.read(8000)  # 先頭8KB程度読んで判定
        
        if not raw:
            return True, 'utf-8'  # 空ファイルはテキスト扱い
        
        result = chardet.detect(raw)
        encoding = result.get('encoding')
        confidence = result.get('confidence', 0)
        
        # 信頼度が低い場合や検出できない場合はバイナリ扱い
        if not encoding or confidence < 0.5:
            return False, None
        
        # 実際に読めるか確認
        try:
            raw.decode(encoding)
            return True, encoding
        except (UnicodeDecodeError, LookupError):
            return False, None
            
    except Exception:
        return False, None


def process_directory(current_path, root_path, output_dir, args, report_items, merger, processed_archives=None, password_protected_files=None):
    """
    ディレクトリを再帰的に処理する関数
    merger: MergedOutputManager instance (None if merge disabled)
    password_protected_files: パスワード保護されたファイルのリスト
    """
    if processed_archives is None:
        processed_archives = set()
    if password_protected_files is None:
        password_protected_files = []

    # current_pathがアーカイブファイルの場合
    if current_path.is_file():
        ext = current_path.suffix.lower()
        if ext in ARCHIVE_EXTENSIONS:
            if current_path in processed_archives:
                return password_protected_files
            processed_archives.add(current_path)
            
            print(f"Extracting Archive [{ext}]: {current_path.name} ...")
            try:
                with tempfile.TemporaryDirectory() as temp_dir:
                    # 圧縮形式に応じて展開
                    result = "OK"
                    if ext == '.zip':
                        result = extract_zip_with_encoding(current_path, temp_dir)
                    elif ext == '.7z':
                        result = extract_7z(current_path, temp_dir)
                    elif ext == '.rar':
                        result = extract_rar(current_path, temp_dir)
                    elif ext in ['.tar', '.gz', '.tgz']:
                        result = extract_tar(current_path, temp_dir)
                    elif ext == '.lzh':
                        result = extract_lzh(current_path, temp_dir)
                    
                    if result == "PASSWORD_PROTECTED":
                        print(f"    [!] Password protected: {current_path.name}")
                        password_protected_files.append(str(current_path))
                    elif result == "OK":
                        process_directory(Path(temp_dir), Path(temp_dir), output_dir, args, report_items, merger, processed_archives, password_protected_files)
            except Exception as e:
                print(f"Error processing archive {current_path}: {e}")
            return password_protected_files

    # ディレクトリ処理
    for root, dirs, files in os.walk(current_path):
        if OUTPUT_DIR_NAME in root or "converted_files_merged" in root:
            continue
            
        for file in files:
            file_path = Path(root) / file
            
            # 隠しファイルをスキップ
            if file.startswith('.'):
                continue
            
            # シンボリックリンクをスキップ
            if file_path.is_symlink():
                print(f"[Skipped Symlink] {file}")
                continue
            
            # 巨大ファイルをスキップ
            try:
                file_size = file_path.stat().st_size
                if file_size > MAX_FILE_SIZE:
                    size_mb = file_size / (1024 * 1024)
                    print(f"[Skipped Large File] {file} ({size_mb:.1f}MB > 100MB)")
                    continue
            except OSError:
                pass  # ファイルサイズ取得失敗は無視
            
            ext = file_path.suffix.lower()
            
            # スキップ対象（明示的にスキップ）
            if ext in SKIP_EXTENSIONS:
                print(f"[Skipped Unsupported] {file}")
                continue
            
            # アーカイブファイルの再帰処理
            if ext in ARCHIVE_EXTENSIONS:
                if file_path not in processed_archives:
                    processed_archives.add(file_path)
                    print(f"Extracting Archive [{ext}]: {file} ...")
                    try:
                        with tempfile.TemporaryDirectory() as temp_dir:
                            result = "OK"
                            if ext == '.zip':
                                result = extract_zip_with_encoding(file_path, temp_dir)
                            elif ext == '.7z':
                                result = extract_7z(file_path, temp_dir)
                            elif ext == '.rar':
                                result = extract_rar(file_path, temp_dir)
                            elif ext in ['.tar', '.gz', '.tgz']:
                                result = extract_tar(file_path, temp_dir)
                            elif ext == '.lzh':
                                result = extract_lzh(file_path, temp_dir)
                            
                            if result == "PASSWORD_PROTECTED":
                                print(f"    [!] Password protected: {file}")
                                password_protected_files.append(str(file_path))
                            elif result == "OK":
                                process_directory(Path(temp_dir), Path(temp_dir), output_dir, args, report_items, merger, processed_archives, password_protected_files)
                    except Exception as e:
                        print(f"Error processing archive {file}: {e}")
                continue

            vis_count = 0
            char_count = 0
            markdown_content = ""
            
            # --- File Type Handling ---

            # 1. 新形式Office (.docx, .xlsx, .pptx) - 視覚密度分析
            if ext == '.docx':
                print(f"Processing: {file}")
                vis_count, char_count = analyze_docx(file_path)
            elif ext == '.xlsx':
                print(f"Processing: {file}")
                vis_count, char_count = analyze_xlsx(file_path)
            elif ext == '.pptx':
                if args.skip_ppt:
                    print(f"Skipping PPT: {file}")
                    continue
                print(f"Processing: {file}")
                vis_count, char_count = analyze_pptx(file_path)

            # Check Density if it was analyzed (新形式のみ)
            if ext in OFFICE_EXTENSIONS_NEW:
                ratio = char_count / vis_count if vis_count > 0 else 9999
                is_dense_visual = ratio < TEXT_PER_VISUAL_THRESHOLD
                if is_dense_visual or vis_count >= 5:
                    print(f"  [Auto-Switch] High density detected (Visuals: {vis_count}). Converting to PDF...")
                    
                    # Target Filename (Flat)
                    target_pdf_name = get_output_filename(root_path, file_path, extension=".pdf")
                    final_pdf_path = output_dir / target_pdf_name
                    
                    # Try PDF Conversion
                    pdf_result = convert_to_pdf_via_libreoffice(file_path, output_dir)
                    
                    if pdf_result:
                        try:
                            if pdf_result != final_pdf_path:
                                if final_pdf_path.exists():
                                    final_pdf_path.unlink()
                                pdf_result.rename(final_pdf_path)
                            
                            report_items.append((file, vis_count, char_count, ratio, f"Converted to PDF"))
                            print(f"    -> Success: {target_pdf_name}")
                            
                            if merger:
                                try:
                                    shutil.copy2(final_pdf_path, merger.output_dir / target_pdf_name)
                                except Exception as e:
                                    print(f"Error copying converted PDF to merged dir: {e}")
                        except Exception as e:
                            print(f"    Error renaming PDF: {e}")
                    else:
                        print("    [Fallback] PDF conversion failed. Copying original file.")
                        report_items.append((file, vis_count, char_count, ratio, "Kept Original (PDF Fail)"))
                        
                        orig_out_name = get_output_filename(root_path, file_path, extension=ext)
                        try:
                            shutil.copy2(file_path, output_dir / orig_out_name)
                        except Exception:
                            pass
                        if merger:
                            try:
                                shutil.copy2(file_path, merger.output_dir / orig_out_name)
                            except Exception:
                                pass
                    continue

                else:
                    # Low Density -> Markdown変換
                    markdown_content = convert_with_markitdown(file_path)

            # 2. PDF Files (Direct Copy)
            elif ext == '.pdf':
                print(f"Copying PDF: {file}")
                output_filename = get_output_filename(root_path, file_path, extension=".pdf")
                
                try:
                    shutil.copy2(file_path, output_dir / output_filename)
                except Exception:
                    pass

                if merger:
                    try:
                        shutil.copy2(file_path, merger.output_dir / output_filename)
                    except Exception as e:
                        print(f"Error copying PDF to merged dir: {e}")
                continue

            # 3. Legacy Office Files (.doc, .xls, .ppt)
            elif ext in OFFICE_EXTENSIONS_LEGACY:
                if ext == '.ppt' and args.skip_ppt:
                    print(f"Skipping PPT (Legacy): {file}")
                    continue
                print(f"Processing Legacy Office[{ext}]: {file}")
                markdown_content = convert_with_markitdown(file_path)
                if markdown_content is None:
                    print(f"    [Warning] Could not convert legacy file: {file}")
                    continue

            # 4. MarkItDown対応形式 (RTF, EPUB, Outlook)
            elif ext in MARKITDOWN_EXTENSIONS:
                print(f"Processing MarkItDown[{ext}]: {file}")
                markdown_content = convert_with_markitdown(file_path)
                if markdown_content is None:
                    print(f"    [Warning] Could not convert file: {file}")
                    continue

            # 5. Visio (.vsdx, .vsd) -> PDF変換
            elif ext in VISIO_EXTENSIONS:
                print(f"Processing Visio: {file}")
                target_pdf_name = get_output_filename(root_path, file_path, extension=".pdf")
                final_pdf_path = output_dir / target_pdf_name
                
                pdf_result = convert_to_pdf_via_libreoffice(file_path, output_dir)
                if pdf_result:
                    try:
                        if pdf_result != final_pdf_path:
                            if final_pdf_path.exists():
                                final_pdf_path.unlink()
                            pdf_result.rename(final_pdf_path)
                        print(f"    -> Success: {target_pdf_name}")
                        if merger:
                            try:
                                shutil.copy2(final_pdf_path, merger.output_dir / target_pdf_name)
                            except Exception as e:
                                print(f"Error copying Visio PDF to merged dir: {e}")
                    except Exception as e:
                        print(f"    Error renaming Visio PDF: {e}")
                else:
                    print(f"    [Warning] Could not convert Visio file: {file}")
                continue

            # 6. 画像 -> PDF変換
            elif ext in IMAGE_EXTENSIONS:
                print(f"Processing Image: {file}")
                target_pdf_name = get_output_filename(root_path, file_path, extension=".pdf")
                final_pdf_path = output_dir / target_pdf_name
                
                pdf_result = convert_image_to_pdf(file_path, output_dir)
                if pdf_result:
                    try:
                        if pdf_result != final_pdf_path:
                            if final_pdf_path.exists():
                                final_pdf_path.unlink()
                            pdf_result.rename(final_pdf_path)
                        print(f"    -> Success: {target_pdf_name}")
                        if merger:
                            try:
                                shutil.copy2(final_pdf_path, merger.output_dir / target_pdf_name)
                            except Exception as e:
                                print(f"Error copying image PDF to merged dir: {e}")
                    except Exception as e:
                        print(f"    Error renaming image PDF: {e}")
                else:
                    print(f"    [Warning] Could not convert image: {file}")
                continue

            # 7. Universal Text Loader (テキストファイル)
            elif ext not in OFFICE_EXTENSIONS_ALL and ext != '.pdf':
                # MIMEタイプでテキスト判定（利用可能な場合）
                mime_is_text = is_likely_text_by_mime(file_path)
                is_known_text = ext in TEXT_EXTENSIONS
                detected_encoding = None
                
                if is_known_text or mime_is_text == True:
                    # テキストファイルとして処理
                    is_readable, detected_encoding = is_text_file(file_path)
                    if not detected_encoding:
                        detected_encoding = 'utf-8'
                elif mime_is_text == False:
                    # MIMEタイプでバイナリと判定
                    print(f"[Skipped Binary] {file}")
                    continue
                else:
                    # MIMEタイプ不明、chardetで判定
                    is_readable, detected_encoding = is_text_file(file_path)
                    if not is_readable:
                        print(f"[Skipped Binary] {file}")
                        continue
                
                print(f"Processing Text[{ext}] ({detected_encoding}): {file}")
                try:
                    with open(file_path, 'r', encoding=detected_encoding, errors='replace') as f:
                        markdown_content = f.read()
                        if not markdown_content.strip():
                            markdown_content = "(Empty File)"
                except Exception as e:
                    print(f"Error reading text file {file}: {e}")
                    markdown_content = ""

            # --- Post-Processing (Report & Write) ---

            # Report Logic (For converted files that might still have some visuals)
            if vis_count > 0 and ext in ['.docx', '.xlsx', '.pptx']:
                 # If we are here, it wasn't high density enough to switch, but maybe worth noting?
                 # Actually, let's strictly log the ones we converted too if they had visuals
                 ratio = char_count / vis_count if vis_count > 0 else 9999
                 report_items.append((file, vis_count, char_count, ratio, "Converted"))

            if markdown_content:
                # 構造を維持したファイル名を生成
                output_filename = get_output_filename(root_path, file_path, extension=".md")
                output_path = output_dir / output_filename
                
                try:
                    rel_path_str = str(file_path.relative_to(root_path))
                except ValueError:
                    rel_path_str = file if isinstance(file, str) else file.name

                metadata_header = f"""# File Info
- Original Filename: {file}
- Relative Path: {rel_path_str}
- Context: {rel_path_str.replace(os.sep, ' > ')}

---
"""
                final_content = metadata_header + markdown_content + "\n\n---\n\n"

                # 1. Write individual file (Always)
                try:
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(final_content)
                except Exception as e:
                    print(f"Failed to write individual {output_path}: {e}")
                
                # 2. Add to Smart Merger (Only if enabled)
                if merger:
                    merger.add_content(output_filename, final_content)
    
    return password_protected_files

def main():
    """
    メインエントリーポイント
    """
    args = setup_args()
    target_path = Path(args.target_dir)
    
    if not target_path.exists():
        print(f"Error: Path '{target_path}' not found.")
        return 1

    # 出力ディレクトリ設定
    if target_path.is_dir():
        output_dir = target_path / OUTPUT_DIR_NAME
        root_processing_path = target_path
    else:
        output_dir = target_path.parent / OUTPUT_DIR_NAME
        root_processing_path = target_path.parent
        
    output_dir.mkdir(exist_ok=True)

    # ログ設定
    logger = setup_logging(output_dir, verbose=args.verbose)
    
    if args.quiet:
        # コンソールハンドラを無効化
        for handler in logger.handlers:
            if isinstance(handler, logging.StreamHandler) and not isinstance(handler, logging.FileHandler):
                handler.setLevel(logging.CRITICAL)
    
    # 処理サマリー初期化
    summary = ProcessingSummary(target_path=str(target_path))
    
    # Dry-run モード
    if args.dry_run:
        logger.info("=== DRY-RUN MODE ===")
        logger.info("Following files would be processed:")
        # TODO: ファイル一覧を表示する処理を追加
        logger.info("Dry-run complete. No files were actually processed.")
        return 0

    # Merged Output Setup
    merger = None
    if args.merge:
        if target_path.is_dir():
            merged_dir = target_path / (OUTPUT_DIR_NAME + "_merged")
        else:
            merged_dir = target_path.parent / (OUTPUT_DIR_NAME + "_merged")
        
        logger.info(f"Target: {target_path}")
        logger.info(f"Output: {output_dir}")
        logger.info(f"Merged: {merged_dir}")
        logger.info("-" * 50)
        
        merger = MergedOutputManager(merged_dir)
    else:
        logger.info(f"Target: {target_path}")
        logger.info(f"Output: {output_dir}")
        logger.info("(Use --merge to create combined output)")
        logger.info("-" * 50)
    
    report_items = []
    password_protected_files = []
    
    password_protected_files = process_directory(
        target_path, root_processing_path, output_dir, args, 
        report_items, merger, password_protected_files=password_protected_files
    )
    
    # Finalize Merge
    if merger:
        merger.finalize()

    # サマリー保存
    summary_file = summary.save(output_dir)

    logger.info("")
    logger.info("=" * 60)
    logger.info(" COMPLETED")
    logger.info("=" * 60)
    
    # 統計サマリー
    logger.info(f"\nProcessing Summary:")
    logger.info(f"  Total files:        {summary.total_files}")
    logger.info(f"  Processed:          {summary.processed}")
    logger.info(f"  Skipped:            {summary.skipped}")
    logger.info(f"  Errors:             {summary.errors}")
    logger.info(f"  Password protected: {summary.password_protected}")
    logger.info(f"\nReport saved to: {summary_file}")
    
    # パスワード保護ファイルレポート
    if password_protected_files:
        logger.warning("\n[!] PASSWORD PROTECTED FILES (Could not process)")
        logger.warning("-" * 60)
        for pf in password_protected_files:
            logger.warning(f"  - {pf}")
        logger.warning("-" * 60)
        logger.warning(f"  Total: {len(password_protected_files)} file(s)")
    
    if report_items:
        logger.info("\n[!] PROCESSED FILE REPORT (Visual Density)")
        logger.info(f" {'Filename':<40} | {'Visuals':<7} | {'Density':<7} | {'Status'}")
        logger.info("-" * 90)
        for (fname, v, c, r, status) in report_items:
             rating = "High" if r < TEXT_PER_VISUAL_THRESHOLD else "Low"
             logger.info(f" {fname:<40} | {v:>7} | {int(r):>5} ({rating}) | {status}")
        logger.info("-" * 90)
        logger.info(" * 'Converted to PDF': High visual density -> Auto-converted to PDF via LibreOffice.")
        logger.info(" * 'Kept Original (PDF Fail)': PDF Conversion failed -> Copied original file.")
    
    return 0


if __name__ == "__main__":
    exit(main())

