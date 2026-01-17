# notebooklm_loader/converters/office_converter.py
"""Office変換モジュール"""

import logging
import docx
import openpyxl
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from markitdown import MarkItDown
from typing import Tuple, Optional


def analyze_docx(file_path) -> Tuple[int, int]:
    """
    Wordファイルの視覚要素と文字数を分析する
    
    Args:
        file_path: 対象ファイルのパス
        
    Returns:
        (visual_count, char_count): 視覚要素数と文字数のタプル
    """
    logger = logging.getLogger("notebooklm_loader")
    try:
        doc = docx.Document(file_path)
        visual_count = 0
        char_count = 0
        for para in doc.paragraphs:
            text = para.text.strip()
            char_count += len(text)
            for run in para.runs:
                if run.element.xpath('.//a:blip'):
                     visual_count += 1
        return visual_count, char_count
    except Exception as e:
        logger.debug(f"analyze_docx error for {file_path}: {e}")
        return 0, 0


def analyze_xlsx(file_path) -> Tuple[int, int]:
    """
    Excelファイルの視覚要素と文字数を分析する
    
    Args:
        file_path: 対象ファイルのパス
        
    Returns:
        (visual_count, char_count): 視覚要素数と文字数のタプル
    """
    logger = logging.getLogger("notebooklm_loader")
    try:
        visual_count = 0
        char_count = 0
        try:
            wb_obj = openpyxl.load_workbook(file_path, data_only=True)
            for sheet_name in wb_obj.sheetnames:
                sheet = wb_obj[sheet_name]
                if hasattr(sheet, '_charts') and sheet._charts:
                    visual_count += len(sheet._charts)
        except Exception as e:
            logger.debug(f"analyze_xlsx chart detection error: {e}")
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            try:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                if not df.empty:
                    csv_text = df.to_csv(index=False)
                    char_count += len(csv_text)
            except Exception as e:
                logger.debug(f"analyze_xlsx sheet {sheet_name} error: {e}")
        return visual_count, char_count
    except Exception as e:
        logger.debug(f"analyze_xlsx error for {file_path}: {e}")
        return 0, 0


def analyze_pptx(file_path) -> Tuple[int, int]:
    """
    PowerPointファイルの視覚要素と文字数を分析する
    
    Args:
        file_path: 対象ファイルのパス
        
    Returns:
        (visual_count, char_count): 視覚要素数と文字数のタプル
    """
    logger = logging.getLogger("notebooklm_loader")
    try:
        prs = Presentation(file_path)
        visual_count = 0
        char_count = 0
        for slide in prs.slides:
            if slide.shapes.title:
                char_count += len(slide.shapes.title.text.strip())
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    char_count += len(shape.text_frame.text.strip())
                is_visual = False
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: is_visual = True
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP: is_visual = True
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                     if not shape.has_text_frame or not shape.text_frame.text.strip(): is_visual = True
                if is_visual: visual_count += 1
        return visual_count, char_count
    except Exception as e:
        logger.debug(f"analyze_pptx error for {file_path}: {e}")
        return 0, 0


def convert_with_markitdown(file_path, max_retries: int = 3) -> Optional[str]:
    """
    MarkItDownを使用してファイルをMarkdownに変換
    
    Args:
        file_path: 対象ファイルのパス
        max_retries: 最大リトライ回数（デフォルト: 3）
        
    Returns:
        変換後のMarkdown文字列、失敗時はNone
    """
    import time
    logger = logging.getLogger("notebooklm_loader")
    
    last_error = None
    for attempt in range(max_retries):
        try:
            md = MarkItDown()
            result = md.convert(str(file_path)) 
            if result and result.text_content:
                return result.text_content
            return ""
        except Exception as e:
            last_error = e
            if attempt < max_retries - 1:
                wait_time = 2 ** attempt  # 指数バックオフ: 1, 2, 4秒
                logger.debug(f"Retry {attempt + 1}/{max_retries} for {file_path.name} after {wait_time}s: {e}")
                time.sleep(wait_time)
            else:
                logger.warning(f"    Error converting {file_path.name} after {max_retries} attempts: {e}")
    
    return None

